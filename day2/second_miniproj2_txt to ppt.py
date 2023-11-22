import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd

def create_pptx(title, content):
    # Create a new PowerPoint presentation
    prs = Presentation()

    # Create title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title

    # Create content slides
    for section in content:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.shapes.placeholders[1]
        title, body = section.split("\n", 1)
        title_shape.text = title
        body_shape.text = body

    # Create thank you slide
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Thank You!"

    # Save the presentation
    filepath = "presentation.pptx"
    prs.save(filepath)
    return filepath

st.title("텍스트를 ppt 파일로 변환하기")

title = st.text_input("제목을 입력하세요")
content = st.text_area("본문을 입력하세요 (각 섹션을 개행으로 구분하세요)", height = 400)
content_sections = content.split("\n\n")

sections = [title] + [sec.split('\n', 1)[0] + "\n" + sec.split('\n', 1)[-1] for sec in content_sections]
display = pd.DataFrame(sections, columns=['Content'])
display.index = [f'슬라이드{num}' for num in range(1, len(content_sections) + 2)]

st.write(display)



if st.button("위와 같이 PPTX 생성할까요?"):
    filepath = create_pptx(title, content_sections)
    st.success(f"PPTX 파일이 생성되었습니다: {filepath}")
